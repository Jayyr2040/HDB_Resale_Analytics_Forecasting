#!/usr/bin/env python3
"""
HDB Analytics Executive Presentation Generator
Generates a professional PowerPoint presentation (9 slides, 10 minutes)
"""

import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("ERROR: python-pptx not installed")
    print("Install with: pip install python-pptx")
    sys.exit(1)

# Color scheme
DARK_BLUE = RGBColor(31, 78, 121)
LIGHT_BLUE = RGBColor(79, 129, 189)
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(68, 68, 68)

def create_presentation():
    """Create the executive presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide1.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BLUE
    
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "HDB RESALE ANALYTICS PLATFORM"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Real-Time Market Intelligence | Automated Pipeline"
    p.font.size = Pt(28)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Executive Summary
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide2.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    
    header = slide2.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE
    
    title_box = slide2.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Executive Summary"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    content_box = slide2.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    items = [
        "🎯 Problem: Manual analysis creates 4-week delays in market insights",
        "",
        "💡 Solution: Automated cloud pipeline with real-time data",
        "",
        "📊 Key Results:",
        "   • 75% faster insights (4 weeks → 5 days)",
        "   • 100% HDB market coverage (230K+ transactions)",
        "   • 258% Year 1 ROI | $96K annual savings",
        "   • Zero manual touchpoints"
    ]
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    # Slide 3: Business Value
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide3.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    
    header = slide3.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE
    
    title_box = slide3.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Business Value: Before vs After"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column
    left_header = slide3.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4), Inches(0.4))
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    p.text = "Before"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    left_items = ["⏱ 4 weeks to insights", "💰 $8K/month labor", "📊 60% data coverage", "🔄 Manual processing", "❌ Stale analysis"]
    left_box = slide3.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.3), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    # Right column
    right_header = slide3.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4), Inches(0.4))
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    p.text = "After"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    right_items = ["⚡ 5 days to insights (5x faster)", "💰 $300/month cloud cost", "📊 100% data coverage", "🤖 Fully automated", "✅ Real-time insights"]
    right_box = slide3.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    # Slide 4: How It Works
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide4.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    
    header = slide4.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE
    
    title_box = slide4.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "How It Works: 3-Layer Pipeline"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    content_box = slide4.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    items = [
        "Layer 1: INGESTION (Dagster)",
        "   Fetch HDB data from data.gov.sg + enrich with geospatial features (20 min)",
        "",
        "Layer 2: TRANSFORMATION (dbt)",
        "   Clean data, validate quality, build star schema for analytics (5 min)",
        "",
        "Layer 3: INSIGHTS (Streamlit)",
        "   Interactive dashboard with executive KPIs, trends, geospatial analysis"
    ]
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # Slide 5: Key Metrics
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide5.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    
    header = slide5.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE
    
    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Key Performance Metrics"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column - Operational
    left_header = slide5.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4), Inches(0.4))
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    p.text = "Operational"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    left_items = ["✅ 230K rows/month processed", "✅ 99.9% pipeline uptime", "✅ <2 sec dashboard load", "✅ 99.2% data quality pass", "✅ Zero manual touchpoints"]
    left_box = slide5.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.3), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    # Right column - Financial
    right_header = slide5.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4), Inches(0.4))
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    p.text = "Financial"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    right_items = ["💵 258% Year 1 ROI", "💵 4.6 month payback", "💵 $300/month (vs $8K)", "💵 160 hours/year freed", "💵 3,200%+ Year 2+ ROI"]
    right_box = slide5.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(15)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(6)
        p.space_after = Pt(6)
    
    # Slide 6: Risk Management
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide6.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    
    header = slide6.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE
    
    title_box = slide6.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Risk Management"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    content_box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    items = [
        "✓ Data Outage → Local cache activates automatically",
        "",
        "✓ API Rate Limits → 30-day intelligent caching (99% hit rate)",
        "",
        "✓ Data Quality Issues → Automated dbt tests flag anomalies",
        "",
        "✓ Pipeline Failures → Dagster alerts + automated recovery",
        "",
        "✓ Compliance & Audit → Version control + BigQuery logs"
    ]
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(4)
        p.space_after = Pt(4)
    
    # Slide 7: Roadmap
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide7.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    
    header = slide7.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE
    
    title_box = slide7.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Roadmap"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Left column
    left_header = slide7.shapes.add_textbox(Inches(0.5), Inches(1.1), Inches(4), Inches(0.4))
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    p.text = "Completed ✅"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    left_items = ["✓ Automated pipeline", "✓ Executive dashboard", "✓ Geospatial analytics", "✓ 96% cost optimization"]
    left_box = slide7.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(4.3), Inches(5.5))
    tf = left_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    # Right column
    right_header = slide7.shapes.add_textbox(Inches(5.2), Inches(1.1), Inches(4), Inches(0.4))
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    p.text = "Next 12 Months 🎯"
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = DARK_BLUE
    
    right_items = ["🎯 Predictive pricing (ML)", "🎯 Real-time alerts", "🎯 Portfolio benchmarking", "🎯 Mobile app"]
    right_box = slide7.shapes.add_textbox(Inches(5.2), Inches(1.6), Inches(4.3), Inches(5.5))
    tf = right_box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)
        p.space_after = Pt(8)
    
    # Slide 8: Q&A
    slide8 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide8.background.fill
    bg.solid()
    bg.fore_color.rgb = WHITE
    
    header = slide8.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(0.8))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_BLUE
    header.line.color.rgb = DARK_BLUE
    
    title_box = slide8.shapes.add_textbox(Inches(0.5), Inches(0.15), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Common Questions"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    content_box = slide8.shapes.add_textbox(Inches(0.8), Inches(1.1), Inches(8.4), Inches(6))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    items = [
        "Q: What if the pipeline breaks?",
        "A: Dagster alerts within minutes + local cache provides continuity",
        "",
        "Q: Is our data secure?",
        "A: Google Cloud SOC 2 certified + full audit trail",
        "",
        "Q: Can non-technical people use this?",
        "A: Yes - zero SQL required, interactive filters, 15-min training",
        "",
        "Q: How do we validate data correctness?",
        "A: dbt tests + quality dashboard + monthly source validation"
    ]
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(3)
        p.space_after = Pt(3)
    
    # Slide 9: Closing
    slide9 = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide9.background.fill
    bg.solid()
    bg.fore_color.rgb = DARK_BLUE
    
    title_box = slide9.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "Ready to Scale"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    subtitle_box = slide9.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(2.5))
    tf = subtitle_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "✅ Delivered on Time & Budget\n✅ Exceeding Performance Metrics\n\nNext Step: Your Approval + Resources to Expand"
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_BLUE
    p.alignment = PP_ALIGN.CENTER
    p.space_before = Pt(12)
    p.space_after = Pt(12)
    
    return prs

if __name__ == "__main__":
    try:
        print("Generating HDB Analytics Executive Presentation...")
        prs = create_presentation()
        
        output_file = "HDB_Analytics_Executive_Presentation.pptx"
        prs.save(output_file)
        
        print(f"✅ SUCCESS!")
        print(f"📁 File: {output_file}")
        print(f"📊 Slides: {len(prs.slides)}")
        print(f"⏱️ Duration: ~10 minutes")
        print(f"\n🎯 Ready to present!")
        
    except Exception as e:
        print(f"❌ Error: {e}")
        sys.exit(1)